import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from jsonschema import Draft202012Validator
from pptx import Presentation


DEFAULT_SCHEMA_PATH = Path("schema/ingested_document.schema.json")


def load_json(path: Path) -> dict[str, Any]:
    """ Load JSON data from file and return as a dict"""
    
    with path.open("r", encoding="utf-8") as file:
        value = json.load(file)

    if not isinstance(value, dict):
        raise TypeError(f"Expected a JSON object in {path}")

    return value


def serialize_datetime(value: datetime | None) -> str | None:
    """ Convert datetime object into a string """
    
    if value is None:
        return None

    return value.isoformat()


def make_document_id(path: Path) -> str:
    """ Generate normalized document ID from file name"""
    
    normalized = re.sub(
        r"[^a-zA-Z0-9]+",
        "_",
        path.stem,
    ).strip("_")

    return normalized.lower() or "presentation"


def normalize_text(value: str) -> str:
    return " ".join(value.split())


def normalize_table(shape: Any) -> tuple[str, int, int]:
    """ Convert PPTX table shape into normalized text representation"""
    
    table = shape.table
    rows: list[str] = []

    for row in table.rows:
        cells = [
            normalize_text(cell.text)
            for cell in row.cells
        ]
        rows.append(" | ".join(cells))

    return (
        "\n".join(rows),
        len(table.rows),
        len(table.columns),
    )


def get_slide_title(slide: Any) -> str | None:
    """ Return normalized title text from a slide"""
    
    title_shape = slide.shapes.title

    if title_shape is None:
        return None

    text = normalize_text(title_shape.text)

    return text or None


def classify_text_shape(
    shape: Any,
    slide_title: str | None,
) -> str:
    """ Identify the type of text shape based on style and content"""
    
    text = normalize_text(shape.text)

    if slide_title and text == slide_title:
        return "title"

    return "paragraph"


def ingest_pptx(source_path: Path) -> dict[str, Any]:
    """ Extract text and tables from a PPTX file """
    
    if source_path.suffix.lower() != ".pptx":
        raise ValueError(
            f"Expected a .pptx file, got: {source_path.suffix}"
        )

    presentation = Presentation(source_path)
    properties = presentation.core_properties

    blocks: list[dict[str, Any]] = []
    plain_text_parts: list[str] = []
    warnings: list[str] = []

    block_index = 0
    table_index = 0

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_title = get_slide_title(slide)

        plain_text_parts.append(
            f"[Slide {slide_number}]"
        )

        # Position-based sorting gives a more natural reading order instead of
        # relying only on the underlying XML shape order.
        shapes = sorted(
            slide.shapes,
            key=lambda shape: (
                getattr(shape, "top", 0),
                getattr(shape, "left", 0),
            ),
        )

        for shape in shapes:
            if getattr(shape, "has_table", False):
                text, row_count, column_count = (
                    normalize_table(shape)
                )

                if not text.strip():
                    warnings.append(
                        "Skipped empty table on "
                        f"slide {slide_number}, "
                        f"table index {table_index}."
                    )
                    table_index += 1
                    continue

                blocks.append(
                    {
                        "block_id": f"block_{block_index}",
                        "block_type": "table",
                        "text": text,
                        "location": {
                            "block_index": block_index,
                            "page": None,
                            "slide": slide_number,
                            "table_index": table_index,
                        },
                        "metadata": {
                            "style": None,
                            "heading_level": None,
                            "rows": row_count,
                            "columns": column_count,
                        },
                    }
                )

                plain_text_parts.append(
                    f"[Table {table_index + 1}]\n{text}"
                )

                block_index += 1
                table_index += 1
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            text = normalize_text(shape.text)

            if not text:
                continue

            block_type = classify_text_shape(
                shape=shape,
                slide_title=slide_title,
            )

            blocks.append(
                {
                    "block_id": f"block_{block_index}",
                    "block_type": block_type,
                    "text": text,
                    "location": {
                        "block_index": block_index,
                        "page": None,
                        "slide": slide_number,
                        "table_index": None,
                    },
                    "metadata": {
                        "style": None,
                        "heading_level": None,
                        "rows": None,
                        "columns": None,
                    },
                }
            )

            plain_text_parts.append(text)
            block_index += 1

        try:
            notes_text = normalize_text(
                slide.notes_slide.notes_text_frame.text
            )
        except Exception:
            notes_text = ""

        if notes_text:
            warnings.append(
                f"Speaker notes detected on slide {slide_number} "
                "but not yet included as ingestion blocks."
            )

    if not blocks:
        warnings.append(
            "No text or tables were extracted from the presentation."
        )

    title = (
        properties.title.strip()
        if properties.title
        else None
    )

    if title is None:
        for block in blocks:
            if block["block_type"] == "title":
                title = block["text"]
                break

    return {
        "document_id": make_document_id(source_path),
        "file_name": source_path.name,
        "file_type": "pptx",
        "title": title,
        "metadata": {
            "author": (
                properties.author.strip()
                if properties.author
                else None
            ),
            "created": serialize_datetime(
                properties.created
            ),
            "modified": serialize_datetime(
                properties.modified
            ),
        },
        "blocks": blocks,
        "plain_text": "\n\n".join(plain_text_parts),
        "warnings": warnings,
    }


def validate_ingested_document(
    document: dict[str, Any],
    schema: dict[str, Any],
) -> None:
    """ Validate ingested document against the JSON schema"""
    
    validator = Draft202012Validator(schema)

    errors = sorted(
        validator.iter_errors(document),
        key=lambda error: list(error.absolute_path),
    )

    if not errors:
        return

    messages: list[str] = []

    for error in errors:
        location = ".".join(
            str(part)
            for part in error.absolute_path
        )
        location = location or "<root>"
        messages.append(f"{location}: {error.message}")

    raise ValueError(
        "Ingested presentation failed schema validation:\n- "
        + "\n- ".join(messages)
    )


def write_ingested_document(
    document: dict[str, Any],
    output_path: Path,
) -> None:
    """ Write the ingested document to a JSON file"""
    
    output_path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    output_path.write_text(
        json.dumps(
            document,
            indent=2,
            ensure_ascii=False,
        )
        + "\n",
        encoding="utf-8",
    )


def parse_arguments() -> argparse.Namespace:
    """ Parse command-line args for PPTX ingestion tool"""
    
    parser = argparse.ArgumentParser(
        description=(
            "Extract slide text and tables from a PPTX file "
            "into the normalized ingestion format."
        )
    )

    parser.add_argument(
        "source",
        type=Path,
        help="Path to the source PPTX file.",
    )
    parser.add_argument(
        "output",
        type=Path,
        help="Path where the ingested JSON will be saved.",
    )
    parser.add_argument(
        "--schema",
        type=Path,
        default=DEFAULT_SCHEMA_PATH,
        help=f"Schema path. Default: {DEFAULT_SCHEMA_PATH}",
    )

    return parser.parse_args()


def main() -> None:
    args = parse_arguments()

    try:
        if not args.source.exists():
            raise FileNotFoundError(
                f"Source file not found: {args.source}"
            )

        schema = load_json(args.schema)
        ingested_document = ingest_pptx(args.source)

        validate_ingested_document(
            document=ingested_document,
            schema=schema,
        )

        write_ingested_document(
            document=ingested_document,
            output_path=args.output,
        )

        print(f"Ingested PPTX: {args.source}")
        print(
            f"Extracted blocks: "
            f"{len(ingested_document['blocks'])}"
        )
        print(
            f"Warnings: "
            f"{len(ingested_document['warnings'])}"
        )
        print(f"Saved normalized document: {args.output}")

    except Exception as error:
        print(
            f"PPTX ingestion failed: {error}",
            file=sys.stderr,
        )
        raise SystemExit(1) from error


if __name__ == "__main__":
    main()